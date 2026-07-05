"""端到端制作工作流编排：deck_outline → 逐页 prompt → SVG → 硬门 → 导出原生 pptx。

串起五层一条链：State(大纲) + Prompt(拼范本) + RAG(范本·在 Prompt 内) + 规则(硬门) + vendor 引擎(导出)。
SVG 生成这步用**依赖注入**（svg_provider）——真实产品 = LLM 调用，测试/手稿 = 传入函数，pipeline 本身可测。
对应 ppt-master 的 plan→svg→pptx，但前置我们的「策略大脑（大纲 + 范本卡 + 硬门）」。fail-closed：大纲非法拒做、硬门 error 可选硬拦。
"""

from __future__ import annotations

import json

import re
from pathlib import Path
from typing import Callable

from engine.chart_template_map import check_image_quota, check_part_visual_quota
from engine.icon_sync import verify_inventory
from engine.prompt_builder import build_deck_prompt
from engine.svg2pptx import export_deck
from reinforce.asset_ledger import build_asset_ledger, merge_generation_records, save_asset_ledger
from reinforce.deck_memory import assert_editable
from reinforce.deck_rules import rules as G
from reinforce.deck_rules.client_tone import run_client_tone_gate
from reinforce.deck_rules.svg_compat import check_svg_compat
from reinforce.deck_rules.visual_review import (check_chart_realized, check_colorblind_safe,
                                                check_icons_realized, check_images_realized,
                                                run_visual_gate)
from reinforce.deck_state import page_input_fingerprint, save_outline, set_page_status, validate_outline
from reinforce.spec_lock import check_against_spec, validate_spec_lock

SvgProvider = Callable[[dict], str]  # ctx(build_deck_prompt 输出) -> svg 字符串


def _extract_text(svg: str) -> str:
    """抽 SVG 所有 <text> 块内文字（含嵌套 tspan）拼接，喂硬门检数字溯源。"""
    blocks = re.findall(r"<text\b[^>]*>(.*?)</text>", svg, re.S)
    return " ".join(re.sub(r"<[^>]+>", "", b).strip() for b in blocks)


def _safe_filename(name: str) -> str:
    """title 直通文件名前的路径字符消毒。

    2026-07-02 saopan扫盘揪出：outline.title 零校验直接拼路径——"FY24/25 增长计划"这类
    中文商业 deck 常见标题会静默建出子目录（页状态已标 done、产物却不在约定位置），
    "../x" 可写出 out_dir。deck_id 有防穿越检查而 title 没有，补齐同款纪律。
    """
    return re.sub(r'[/\\:*?"<>|]', "_", name).strip() or "deck"


def build_deck(outline: dict, svg_provider: SvgProvider, out_dir, *,
               strict: bool = False, notes_map: dict | None = None, spec_lock: dict | None = None,
               render_previews: bool = False, workspace: dict | None = None,
               resume: bool = False) -> dict:
    """跑完整 deck：校验大纲 → 逐页(拼 prompt → 生成 SVG → 硬门) → 导出原生可编辑 pptx。

    svg_provider(ctx)->svg：真实 = Claude 逐页手写（ppt-master 验证过的纪律：禁脚本批量数据驱动模板生成，
    见制作工作流 SKILL）；此处 = 手写 / 注入（测试场景）。
    spec_lock 传入时：① 拼进每页 prompt 强制"重读"（不许凭记忆现编颜色/字体）
    ② 每页生成后机检用色是否越出锁定色板（warn 喂人审，防跨页风格漂）③ 校验 spec_lock 自身字段完整
    （validate_spec_lock，不完整直接拒做·fail-closed）。不传则每页记一条 warn（阶段1.5 是强制关卡，
    见 spec_lock.py；warn 非 error 保留旧调用/单测不依赖 spec_lock 时仍可跑，但不再是完全静默的跳过）。
    页型容量硬门（G.check_page_budget）：page.facets.chart 命中 PAGE_BUDGET 且有 evidence 时用
    len(evidence) 代理条目数机检（2026-07-01 yanjiu研究驱动评估揪出此前零接线，见函数体内注释）。
    strict=True 时任一页 error 即抛（fail-closed 硬）；默认产草稿 + 标记 all_passed=False 供人审。
    render_previews=True 时每页额外渲染一张 `page_N.png`（D7路线A"渲染→自检"闭环用的预览图，
    见 `engine/render_preview.py`）——**缺省 False**：playwright 是可选依赖 + 每次渲染要开关浏览器，
    默认开会拖慢所有调用方（含测试套件）。未装 playwright 时传 True 直接抛 ImportError——
    2026-07-02 saopan扫盘揪出：此前静默跳过（renderer=None 连 warn 都不记），而制作 SKILL 明确
    承诺"未装会抛清楚的 ImportError，不会静默跳过让你误以为自检发生过"——显式要求的自检关卡
    无声蒸发是假绿，fail-closed 该抛就抛（不传 render_previews 的调用方不受影响）。

    workspace（2026-07-02 ④数据产物层接线·依赖注入同 svg_provider 一条纪律，pipeline 不 import
    deck_workspace，只按键消费 dict——deck_workspace.new_workspace 的返回值直接传入即可）：
    ① deck_id + published_ledger 在 → 开工前 assert_editable（**铁律4真接线**：已交付 deck 拒重做，
      此前 mark_published/assert_editable 全仓零生产调用、只读锁形同虚设，同 spec_lock 那次
      "文档喊强制、代码可绕过"是一种病）
    ② outline 键在 → 每页状态流转后落盘 outline.json（活文件·中断后它就是断点续做的入口）
    ③ asset_ledger 键在 → 收尾自动汇集素材/证据溯源台账落盘（见 reinforce/asset_ledger.py）。
    不传 workspace 时返回值 deck_warnings 记一条（deck 级警示一条即可，不逐页刷屏；同 spec_lock
    的 warn 一个道理——允许旧调用/单测跑，但"这份 deck 没有只读锁保护、状态不落盘"不再静默）。

    页状态生命周期（2026-07-02 起真实流转·此前 status 永远停在 planned 是摆设）：
    生成 SVG 后 planned→drafted；硬门全过 drafted→gated；导出成功 gated→done（没过门的页停在
    drafted 供人审返工，不升 done）。resume=True 时：status 已是 gated/done 且 out_dir 下对应
    SVG 还在的页，跳过 svg_provider（省的是 LLM 手写这步贵的），**硬门照样重跑**——规则可能已
    升级，上次过门不代表这次过（防假绿），重检没过如实计入 all_passed/strict。
    返回 {pptx_path, pages:[{n, gate, svg_path, preview_path, resumed}], all_passed, deck_warnings}。
    """
    v = validate_outline(outline)
    if not v["valid"]:
        raise ValueError(f"大纲非法·拒做（fail-closed）：{v['issues']}")
    deck_warnings: list[dict] = []
    ws = workspace or {}
    if ws.get("deck_id") and ws.get("published_ledger"):
        # 铁律4：已交付 deck 只读，重做 = 改已交付产物，PermissionError 硬拦（改需新 deck_id 走新版本）
        assert_editable(ws["deck_id"], ws["published_ledger"])
    if workspace is None:
        deck_warnings.append({"rule": "workspace_missing", "sev": "warn",
            "note": "未传 workspace：本 deck 无已交付只读锁保护、页状态/溯源台账不落盘，"
                    "中断后无法断点续做（④数据产物层约定，见 reinforce/deck_workspace.py）"})
    elif not (ws.get("deck_id") and ws.get("published_ledger")):
        # 2026-07-02 saopan扫盘揪出：手拼残缺 workspace dict（缺 deck_id/published_ledger 键）
        # 此前既不拦也不警——只读锁静默不生效，比"不传"更隐蔽（调用方自以为有保护）。
        deck_warnings.append({"rule": "workspace_incomplete", "sev": "warn",
            "note": "workspace 缺 deck_id/published_ledger 键：已交付只读锁不生效——"
                    "请用 deck_workspace.new_workspace 生成完整 workspace，别手拼 dict"})
    if spec_lock is not None:
        # 2026-07-01 yanjiu研究驱动评估揪出：validate_spec_lock 写好后全仓零调用点(architecture.py
        # 自己标"人工调用一次")——传了不完整的 spec_lock(缺 palette/font_family，或 per_page 声明
        # 一半)此前完全没有机检兜底，直接带着残缺规范滑到生成环节。传了就必须是合法的（fail-closed，
        # 不完整等同没锁定，比"完全不传"更危险——后者至少每页会收到下方 warn 提醒）。
        sv = validate_spec_lock(spec_lock)
        if not sv["valid"]:
            raise ValueError(f"spec_lock 不完整·拒做（fail-closed）：{sv['issues']}")
        # 批④孤岛接线（第5次孤岛复发·2026-07-03）：icon_sync.verify_inventory 写好后全仓零生产
        # 调用点——两份 SKILL 都承诺"icon_sync 存在性验证缺名强制重选（FR7.5）"，但没有任何东西
        # 真跑它：锁定不存在的库/缺名图标一路滑进生成环节，页面出空 <use> 占位。spec_lock 声明
        # icons 即在此逐名验证文件真存在（validate_spec_lock 只查字段形态，不查磁盘）。
        icons = spec_lock.get("icons") or {}
        if icons:
            iv = verify_inventory(icons.get("library"), icons.get("inventory"))
            if not iv["valid"]:
                raise ValueError(
                    f"icons 清单校验失败·拒做（fail-closed·D18 FR7.5 缺名强制重选）："
                    f"library={iv['library']!r}"
                    f"{'（库不存在·可选库见 icon_sync.list_icon_libraries()）' if not iv['library_exists'] else ''}"
                    f" missing={iv['missing']}")
        # 2026-07-02 saopan扫盘揪出（孤岛A1）：CVD 色盲检测写好后 architecture.py 登记
        # "建 spec_lock 时人工跑一次"，但制作 SKILL 阶段1.5 工序通篇无此步骤——没有任何东西
        # 保证它被跑，全仓零生产调用点。接线到这里：spec_lock 传入即对锁定色板跑一次，
        # warn 喂 deck_warnings（色板是 deck 级的，跑一次不逐页刷）。
        deck_warnings.extend(check_colorblind_safe(spec_lock.get("palette") or {}))
        # D18 FR3.5 Part 级视觉配额（warn·deck 级跑一次）：某 Part 全组零强视觉页 → 提示
        # ——对齐真样本"每章 1-2 页强视觉"节奏感；不要求每页有视觉主体（防装饰性噱头）。
        deck_warnings.extend(check_part_visual_quota(outline, spec_lock))
        # D27 素材在场未用机检（第七轮某品牌单确诊·error 拦在逐页生成前）：用户特意投放了品牌
        # 素材（brand_assets 有真文件），spec_lock 的色板来源却不是 brand_vi——第七轮实测：素材
        # 异步投放（开场问"怎么给"→中途放入）而扫描点只有设计定稿一刻，13 页做完才发现素材、
        # 整单返工。本机检把发现时刻从"做完"提前到"开工"（返工成本≈0：改 spec_lock 重跑）。
        # 豁免口：palette_source.assets_reviewed=true = 会话已扫过素材、有理由不用（如素材是
        # 历史 deck 无标准色号）——显式留痕的弃用是决策，静默漏扫是事故，只拦后者。
        _ba = ws.get("brand_assets")
        if _ba and Path(_ba).is_dir():
            _has_assets = any(
                f.is_file() and f.name != "_放什么.md" and not f.name.startswith(".")
                for f in Path(_ba).rglob("*"))
            _ps = spec_lock.get("palette_source") or {}
            if _has_assets and _ps.get("kind") != "brand_vi" and not _ps.get("assets_reviewed"):
                raise ValueError(
                    f"品牌素材在场未用·拒做（fail-closed·D27）：brand_assets/ 有用户投放的素材，"
                    f"但 palette_source.kind={_ps.get('kind')!r} 不是 brand_vi 且无 assets_reviewed "
                    f"豁免——回设计定稿扫描素材做品牌锚定（list_brand_assets→提取色板/字体），"
                    f"或确认弃用后在 palette_source 里加 assets_reviewed=true + 弃用理由（如实留痕）")
        # D22 图像配额分级执法：浓郁档不达标=error 硬拦（直通档下 warn 整批架空[人审
        # waiver]，用户拍板的浓郁承诺必须有 fail-closed 兜底），标准档零图像位 warn。
        img_issues = check_image_quota(outline, spec_lock)
        img_hard = [i for i in img_issues if i.get("sev") == "error"]
        if img_hard:
            raise ValueError(f"视觉丰富度未兑现·拒做（fail-closed·D22）：{img_hard}")
        deck_warnings.extend(img_issues)

    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    svg_paths: list[Path] = []
    page_results: list[dict] = []
    all_passed = True

    renderer = None
    if render_previews:
        from engine.render_preview import SvgRenderer, preview_available
        if not preview_available():
            # 2026-07-02 saopan扫盘揪出：此前 renderer 留 None 静默跳过整个渲染自检关——
            # 显式传了 render_previews=True 却无声蒸发，违背制作 SKILL"不会静默跳过"承诺。
            raise ImportError(
                "render_previews=True 但 playwright/chromium 未装·渲染自检关无法执行"
                "（fail-closed 拒静默跳过）：装 playwright 后重跑，或去掉 render_previews")
        renderer = SvgRenderer().__enter__()  # 全 deck 复用一个浏览器会话，不逐页开关

    try:
        for p in sorted(outline["pages"], key=lambda x: x["n"]):
            n = p["n"]
            sp = out_dir / f"page_{n}.svg"
            resumed = False
            if p.get("status") in ("gated", "done"):
                if resume and sp.is_file():
                    # 断点续做：复用已过门页的 SVG，跳过 svg_provider（LLM 手写是贵的那步）——
                    # 但下方硬门照常重跑（规则可能升级过，上次过门不代表这次过·防假绿）。
                    # 复用前先验输入指纹（D16 🟡②·Bazel action key 思路）：claim/facets/spec_lock
                    # 变了 = SVG 是过期产物，硬门重跑逮不住这种漂移（check_title 只看 claim 参数
                    # 字符串，不校验 SVG 内渲染文字），必须拒复用。
                    fp_now = page_input_fingerprint(p, spec_lock)
                    fp_was = p.get("input_fingerprint")
                    if fp_was and fp_was != fp_now:
                        raise ValueError(
                            f"p{n} 输入已变（claim/facets/spec_lock 与产出该页 SVG 时不一致）·"
                            f"拒复用过期产物：改回原输入，或把该页 status 改回 planned 重做")
                    if not fp_was:
                        # 无指纹的旧页（本机制上线前产的）：无从验证是否过期，warn 不拦——
                        # 硬门仍会重跑兜一部分底，但输入漂移检测对这页是盲的，如实说。
                        deck_warnings.append({"rule": "fingerprint_missing", "sev": "warn",
                            "note": f"p{n} 无输入指纹（旧产物）·复用时无法验证 SVG 是否过期"})
                    svg = sp.read_text(encoding="utf-8")
                    resumed = True
                else:
                    # 已过门的页不许被悄悄重写（同只读锁一个精神·状态是真相源）：
                    # 想复用传 resume=True；确要返工把该页 status 手工改回 planned（罕见路径
                    # 不设机器通道，见 deck_state.STATUS_FLOW 注释）；svg 文件丢了 = 状态与
                    # 产物矛盾，同样值得人来看一眼而不是静默重做。
                    raise ValueError(
                        f"p{n} 已 {p['status']}（硬门已过{'/已导出' if p['status'] == 'done' else ''}）"
                        f"·拒绝静默重做：复用请传 resume=True"
                        f"{'（且 ' + str(sp) + ' 缺失·状态与产物矛盾）' if resume else ''}；"
                        f"确要返工请把该页 status 改回 planned")
            if not resumed:
                ctx = build_deck_prompt(outline, n, spec_lock=spec_lock)  # Prompt 引擎：RAG 范本 + 规范 + 依赖 + spec_lock
                svg = svg_provider(ctx)                   # SVG 生成槽位（真实 = LLM）
                set_page_status(outline, n, "drafted")    # 生命周期：SVG 已生成
                p["input_fingerprint"] = page_input_fingerprint(p, spec_lock)  # 记产出时输入指纹（resume 失效判定用·D16 🟡②）
            # D19 FR5.1+FR1.3 豁免集合：品牌词（"059"曾被当统计数字 error 误报）+ 该页
            # client_provided 证据的数字值（客户自己的资料页面不标源·内部台账已溯源）
            exempt = list(outline.get("brand_terms") or [])
            for _e in (p.get("facets") or {}).get("evidence") or []:
                if _e.get("source_type") == "client_provided":
                    exempt += re.findall(r"\d[\d,.]*%?", str(_e.get("data") or ""))
            gate = G.run_text_gate(_extract_text(svg), title=p.get("claim"), exempt_terms=exempt)  # 硬门
            gate["issues"] = gate["issues"] + check_svg_compat(svg)  # PPT 导出兼容性机检(shared-standards.md)
            gate["issues"] = gate["issues"] + run_visual_gate(svg)  # 视觉缺陷结构级机检(H1/H4/H8/S6·visual-review.md)
            # D18 FR5.2 对客文案门（D19 FR1.2/D20 FR2.1 起分级：brief定位符/字段标签冒号形态/
            # 元叙述组合=error 参与 gate.passed[strict 拦导出]；流程自述/白名单外缩写/工作口语=warn 喂人审）
            gate["issues"] = gate["issues"] + run_client_tone_gate(svg)
            # D18 FR3.2 chart 兑现机检（warn·第一轮实测声明 20 种 chart 兑现 0——渲染后按声明查
            # 期望视觉元素，查不到=可能没兑现喂人审；未知 chart 名不查·宁漏勿噪）
            _declared_chart = (p.get("facets") or {}).get("chart")
            if _declared_chart:
                gate["issues"] = gate["issues"] + check_chart_realized(svg, _declared_chart)
            facets = p.get("facets") or {}
            chart = facets.get("chart")
            evidence = facets.get("evidence") or []
            if chart in G.PAGE_BUDGET and evidence:
                # 2026-07-01 yanjiu研究驱动评估揪出：check_page_budget 此前只有单测覆盖，从未接入
                # 生产链路。count 用 len(evidence) 代理——数据论断页一条 evidence 对应图上一个数据点
                # (donut一瓣/柱状一根/流程一步)，这个字段 to_outline() 已经从 storyline 带下来，
                # 不需要给 outline schema 新增字段。facets.chart 而非 page_function 才是真实图型技术名
                # （page_function 优先取 page_type 兜底，真实 deck 里常是"数据论断"这类不匹配 PAGE_BUDGET 的值）。
                gate["issues"] = gate["issues"] + G.check_page_budget(chart, len(evidence))
            gate["passed"] = gate["passed"] and not any(i["sev"] == "error" for i in gate["issues"])
            if spec_lock:
                gate["issues"] = gate["issues"] + check_against_spec(svg, spec_lock)  # spec_lock 用色机检
            else:
                # 2026-07-01 yanjiu研究驱动评估揪出：spec_lock.py 自己的模块docstring写"不是建议，
                # 是强制关卡"，但此前不传 spec_lock 时 build_deck 完全静默跳过——文档喊"强制"、代码
                # 允许无声绕过，是这几天反复在堵的"文档↔代码不一致"同一种病，只是这次长在机检自己身上。
                # warn 非 error：不破坏现有不传 spec_lock 的调用方/单测（fixture 场景合理不传），
                # 但至少让"这份 deck 没锁定色板/字体"变成 gate 里看得见的信号，不再是彻底沉默。
                gate["issues"] = gate["issues"] + [{"rule": "spec_lock_missing", "sev": "warn",
                    "note": "未传 spec_lock：本页用色/字体未经锁定校验，跨页风格漂移无法机检兜底"
                            "（阶段1.5 强制关卡，见 spec_lock.py 模块说明）"}]
            if not gate["passed"]:
                all_passed = False
                if strict:
                    raise ValueError(f"p{n} 硬门未过（fail-closed）：{gate['issues']}")
                if resumed:
                    # 2026-07-02 saopan扫盘揪出（三路交叉确认）：resumed 页重检失败后状态停在
                    # gated/done，收尾循环把 gated 无条件升 done——落盘 outline.json 假绿，
                    # 门败只活在被丢弃的返回值里、跨会话被遗忘。降级 drafted 落盘如实返工
                    # （STATUS_FLOW 已补 gated/done→drafted 返工边）。
                    set_page_status(outline, n, "drafted")
            elif not resumed:
                set_page_status(outline, n, "gated")      # 生命周期：硬门全过（没过则停在 drafted 供返工）
            if not resumed:
                sp.write_text(svg, encoding="utf-8")
            svg_paths.append(sp)
            preview_path = None
            if renderer is not None and not resumed:
                preview_path = str(renderer.render(svg, out_dir / f"page_{n}.png"))
            elif resumed and (out_dir / f"page_{n}.png").is_file():
                preview_path = str(out_dir / f"page_{n}.png")   # 复用上次的预览图
            page_results.append({"n": n, "gate": gate, "svg_path": str(sp),
                                 "preview_path": preview_path, "resumed": resumed})
            if ws.get("outline"):
                save_outline(outline, ws["outline"])      # 逐页落盘（防中断丢进度·断点续做的入口）
    finally:
        if renderer is not None:
            renderer.__exit__(None, None, None)

    if ws.get("pptx"):
        # 2026-07-02 saopan扫盘揪出：此前成品写 out_dir/<title>.pptx（SKILL 指引 out_dir=ws["pages"]
        # → 成品落 pages/ 下），而 publish_deck 算 sha256 找的是 workspace 约定的 deck.pptx——
        # 两个路径全仓无人搬运，D16 交付物 checksum 机制在标准链路上 100% 空转（永远 no_checksum）。
        # workspace 传入时成品统一落 ws["pptx"]（ARTIFACTS 约定位置），checksum/verify 链路才闭合。
        pptx_path = Path(ws["pptx"])
    else:
        pptx_path = out_dir / f"{_safe_filename(outline['title'])}.pptx"
    if not export_deck(svg_paths, pptx_path):        # vendor 引擎：原生可编辑导出
        raise RuntimeError("vendor 引擎导出失败")
    if notes_map:                                     # 演讲者备注（思路讲解·演讲版「口头补」的脚本）
        _attach_notes(pptx_path, [notes_map.get(pr["n"], "") for pr in page_results])

    for p in outline["pages"]:                        # 生命周期：导出成功 → 过门页升 done
        if p.get("status") == "gated":                # （drafted=硬门没过的页停在原地，如实不升）
            set_page_status(outline, p["n"], "done")
    if ws.get("outline"):
        save_outline(outline, ws["outline"])
    # D26 兑现机检（deck 级·warn 喂人审）：声明了图标清单/图像清单，成品里却整批没兑现——
    # 二轮扫盘确诊 5/6 deck 图标 0 占位、某品牌 9 图 0 落地，此前无任何机检能发现。
    if spec_lock is not None and page_results:
        _svg_texts = [Path(pr["svg_path"]).read_text(encoding="utf-8") for pr in page_results]
        deck_warnings.extend(check_icons_realized(_svg_texts, spec_lock))
        deck_warnings.extend(check_images_realized(_svg_texts, spec_lock,
                                                   images_dir=ws.get("images")))
    if ws.get("asset_ledger"):                        # 素材/证据溯源台账：自动汇集落盘（声明式·不用人另填）
        svgs = {pr["n"]: Path(pr["svg_path"]).read_text(encoding="utf-8") for pr in page_results}
        ledger = build_asset_ledger(outline, svgs=svgs)  # svgs=图像素材段(D16 🟢②)
        # D26 二轮扫盘🔴：重建前先读旧台账，把 AI 生图记录(kind=generation·维权证据链)搬运
        # 保留——此前每次收尾(含续做/重跑)从头重建直接清零 record_image_generation 的留痕。
        old_path = Path(ws["asset_ledger"])
        if old_path.is_file():
            try:
                old = json.loads(old_path.read_text(encoding="utf-8"))
            except (json.JSONDecodeError, UnicodeDecodeError):
                old = None   # 旧台账损坏·如实不搬(重建覆盖)·不让坏文件崩掉收尾
            merge_generation_records(ledger, old)
        save_asset_ledger(ledger, ws["asset_ledger"])

    return {"pptx_path": str(pptx_path), "pages": page_results, "all_passed": all_passed,
            "deck_warnings": deck_warnings}


def _attach_notes(pptx_path, notes_list) -> None:
    """用 python-pptx 给每页写 speaker notes（按 slide 顺序）。notes 是写操作，不触发 3.14 读 bug。"""
    from pptx import Presentation
    pres = Presentation(str(pptx_path))
    for slide, note in zip(pres.slides, notes_list):
        if note:
            slide.notes_slide.notes_text_frame.text = note
    pres.save(str(pptx_path))
