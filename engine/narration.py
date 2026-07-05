"""把讲解词转成机器语音、嵌进 pptx 每页并设自动播放——deck 会自己讲。

链路：讲解词 → TTS → 音频 → python-pptx add_movie 嵌入 → 改 slide timing 让音频翻到该页**自动播放**。
讲解词来源（D18 FR7.8）：`narrate_pptx(..., storyline=sl)` 可直接吃策略层 storyline 推演链
（sowhat/framing/bridge_from/part → `speaker_notes.narration_list_for_pptx`），不再要求调用方
先把"页面文字"拼成列表；LLM 手写讲解稿仍走 narration_list 主路径。
TTS 双引擎（都免费·无 API key）：
  · edge（默认）= 微软 Edge 神经 TTS（edge-tts·中文接近真人·要联网·无 key/无注册）；
  · say        = macOS 自带（本地·无网·机器味·离线兜底）。
对应 ppt-master 的语音讲解。

⚠️ 自动播放的实际效果需在 PowerPoint/Keynote 打开验收（本模块只保证音频嵌入 + timing 改成 delay=0）。
"""

from __future__ import annotations

import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

EDGE_VOICE = "zh-CN-YunxiNeural"  # 云希·男声·神经语音（接近真人）。女声可换 zh-CN-XiaoxiaoNeural
SAY_VOICE = "Tingting"               # macOS 中文女声（机器味）
_MIME = {".mp3": "audio/mpeg", ".m4a": "audio/mp4"}


def tts(text: str, out_audio, engine: str = "edge", voice: str | None = None) -> Path:
    """文本 → 语音文件。engine='edge'（免费神经·无 key·要联网）| 'say'（macOS 本地·无网·机器味）。"""
    out_audio = Path(out_audio)
    if engine == "edge":
        import asyncio
        import edge_tts
        out_mp3 = out_audio.with_suffix(".mp3")

        async def _go():
            await edge_tts.Communicate(text, voice or EDGE_VOICE).save(str(out_mp3))

        asyncio.run(_go())
        return out_mp3
    # macOS say 兜底（本地·无网）
    aiff, m4a = out_audio.with_suffix(".aiff"), out_audio.with_suffix(".m4a")
    try:
        subprocess.run(["say", "-v", voice or SAY_VOICE, "-o", str(aiff), text], check=True)
        subprocess.run(["afconvert", str(aiff), str(m4a), "-f", "m4af", "-d", "aac"], check=True)
    finally:
        # 2026-07-03 二轮扫盘批D：afconvert 失败此前直接抛、.aiff 中间产物泄漏——finally 保证必清
        aiff.unlink(missing_ok=True)
    return m4a


def _set_autoplay(slide) -> None:
    """把该页音频触发条件从「手动(indefinite)」改成「0 延迟自动播放」。"""
    for cond in slide._element.iter(qn("p:cond")):
        if cond.get("delay") == "indefinite":
            cond.set("delay", "0")


def narrate_pptx(pptx_path, narration_list=None, engine: str = "edge",
                 voice: str | None = None, audio_dir=None, *, storyline: dict | None = None) -> dict:
    """给 pptx 每页配语音讲解（自动播放）。engine 见 tts。narration_list 按 slide 顺序，空串=不配音。

    D18 FR7.8（备注叙事）：讲解词来源二选一——
    - narration_list：LLM 按制作 SKILL 阶段 5 备注标准逐页手写的讲解稿（主路径·亮眼）；
    - storyline=：直接传策略层 storyline dict，讲解词从**推演链**生成
      （`speaker_notes.narration_for_storyline`·吃 sowhat/framing/bridge_from/part），
      输入不再是"页面文字的复读"——结构化兜底路径。
    两个都传/都不传直接拒（fail-closed·讲稿来源必须唯一，防"到底在念哪份稿"的歧义）。
    向后兼容：既有调用都是位置参数传 narration_list，行为不变。
    """
    if (narration_list is None) == (storyline is None):
        raise ValueError("narration_list 与 storyline 必须二选一"
                         "（都传/都不传都拒——讲稿来源要唯一，D18 FR7.8）")
    if narration_list is None:
        from engine.speaker_notes import narration_list_for_pptx  # 延迟 import·保持模块轻依赖
        narration_list = narration_list_for_pptx(storyline)
    pptx_path = Path(pptx_path)
    audio_dir = Path(audio_dir or pptx_path.parent / "_audio")
    audio_dir.mkdir(parents=True, exist_ok=True)
    pres = Presentation(str(pptx_path))
    n = 0
    for i, (slide, text) in enumerate(zip(pres.slides, narration_list), start=1):
        if not text:
            continue
        audio = tts(text, audio_dir / f"p{i}", engine=engine, voice=voice)
        slide.shapes.add_movie(str(audio), Inches(12.4), Inches(0.15), Inches(0.45), Inches(0.45),
                               mime_type=_MIME.get(audio.suffix, "audio/mpeg"))
        _set_autoplay(slide)
        n += 1
    pres.save(str(pptx_path))
    return {"pages_narrated": n, "engine": engine}
