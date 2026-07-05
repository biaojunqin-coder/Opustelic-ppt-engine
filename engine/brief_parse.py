"""brief 解析（策略阶段1）——把 brief 文件（PDF/Word/PPT/文本）提取成 LLM-ready 文本。

**免费·本地·无 key**：fitz(PDF) + python-docx(Word) + python-pptx(PPT) + 纯文本。
docling / MinerU 是可选增强（版面分析更强·但包重），默认走轻量提取。对应策略 SKILL 阶段1。

2026-07-01 补·docling 接入（D13 决策记录）：MinerU 因 `requires_python<3.14` 跟本仓库 Python 3.14
硬不兼容（装不上，非网络问题），找替代验证出 docling（MIT 许可·真实测过中文+复杂表格提取，见
`_parse_pdf_with_docling` 文档字符串），`prefer_docling=True` 是当前推荐的增强路径，`prefer_mineru`
留给未来 Python<3.14 场景或 MinerU 补上 3.14 支持后用。

2026-07-01 补·`prefer_mineru` 从空转桩改真接（D6 决策记录🟡走样需修）：此前 `try: import mineru`
成功后什么也不做、原样落回 fitz——比"没做"更容易被误读成"部分接入"。现在真调用
`mineru.cli.common.do_parse`（backend="pipeline"·唯一纯 CPU 可跑+"无幻觉"的后端，契合"宁可漏识别
不可瞎编"）。⚠️ 诚实边界：`do_parse` 是 CLI 导向的批处理 API，需要 MinerU 自己的模型权重
（几百MB~数GB，未装会在首次调用时自动下载），本仓库开发环境未装 `mineru` 包也未下载模型权重，
**这层真接过结构级测试（mock 验证调用参数正确 + 异常/未装两条 fallback 路径），但未跑通过一次
真实模型推理**——跟"研究到位、实装未动"的 waterfall/gantt 是同一类诚实标注，不同的是这里"动"的是
接线代码，不是算法代码。

2026-07-01 补·表格/图片清单/超链接：六路深挖 ppt-master 资产时发现它的 `scripts/source_to_md/`
（doc_to_md.py/ppt_to_md.py 等）提取深度比这里原先的纯文本遍历高一截——pptx 表格内容之前完全
静默丢失（`shape.text_frame.text` 取不到表格）、图片存在与否没有任何记录、超链接也没提取。
这里只挪了"brief 阅读够用"的部分（表格转 markdown 表 / 图片计数提示 / 超链接转内联 `[文字](url)`），
没有挪它的资产管理那层（occurrence 坐标追踪 / 图片落盘 / manifest.json——那是给"文档发布转换"用的，
brief 只是读进来给 LLM 当上下文，不需要那层重量级管理），也没挪它处理"幻灯片内部跳转"的私有 API
hack（`run._r` 那段，brief 阅读用不上，且耦合 python-pptx 内部实现，能不碰就不碰）。
"""

from __future__ import annotations

from pathlib import Path


def _rows_to_md_table(rows: list[list[str]]) -> str:
    """二维字符串表格 → markdown 表格（pptx/docx 表格共用）。"""
    if not rows:
        return ""
    header, *body = rows
    sep = ["---"] * len(header)
    lines = ["| " + " | ".join(header) + " |", "| " + " | ".join(sep) + " |"]
    lines += ["| " + " | ".join(r + [""] * (len(header) - len(r))) + " |" for r in body]
    return "\n".join(lines)


def _clean_cell(text: str) -> str:
    return (text or "").strip().replace("|", "\\|").replace("\n", " ")


def _pptx_run_text(run) -> str:
    """单个 run 转文字；带外部超链接转内联 `[文字](url)`。"""
    text = run.text or ""
    try:
        url = run.hyperlink.address
    except AttributeError:
        url = None
    return f"[{text}]({url})" if (url and text.strip()) else text


def _pptx_table_md(table) -> str:
    """PPT 表格 → markdown 表格（之前 `shape.text_frame.text` 取不到表格内容，整张表静默丢失）。"""
    rows = [[_clean_cell(cell.text) for cell in row.cells] for row in table.rows]
    return _rows_to_md_table(rows)


def _pptx_shape_text(shape) -> str:
    """单个 shape → 文字：表格走 markdown 表，文字框逐段保留超链接。"""
    if shape.has_table:
        return _pptx_table_md(shape.table)
    if shape.has_text_frame and shape.text_frame.text.strip():
        paras = []
        for para in shape.text_frame.paragraphs:
            text = "".join(_pptx_run_text(r) for r in para.runs) or para.text
            if text.strip():
                paras.append(text)
        return "\n".join(paras)
    return ""


def _iter_pptx_shapes(shapes):
    """深度优先展开 GroupShape，产出所有叶子 shape（保持组内 z-order 顺序）。

    2026-07-02 saopan扫盘揪出：旧版只遍历 slide 顶层 shapes——GroupShape 没有 text_frame，
    组合形状里的文字整组静默丢失（真实 brief 里"几个框+连线圈成一组"极常见）。
    shape_type 读取包 try：python-pptx 在 Python 3.14 下个别形状读回会抛（同主循环的兼容处理），
    读不出类型就当叶子交给主循环的 try/except 兜底，不让一个坏形状中断整页提取。
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for sh in shapes:
        try:
            is_group = sh.shape_type == MSO_SHAPE_TYPE.GROUP
        except Exception:
            is_group = False
        if is_group:
            yield from _iter_pptx_shapes(sh.shapes)
        else:
            yield sh


def _docx_para_text(paragraph) -> str:
    """单个段落转文字，按文档原始顺序把超链接转内联 `[文字](url)`（python-docx 1.1+ `iter_inner_content`）。"""
    import docx.text.hyperlink as _h
    parts = []
    for item in paragraph.iter_inner_content():
        if isinstance(item, _h.Hyperlink):
            if item.text.strip():
                parts.append(f"[{item.text}]({item.address})")
        else:
            parts.append(item.text)
    return "".join(parts)


def _docx_table_md(table) -> str:
    rows = [[_clean_cell(cell.text) for cell in row.cells] for row in table.rows]
    return _rows_to_md_table(rows)


def _parse_pdf_with_docling(path: Path) -> str | None:
    """真调用 docling 解析 PDF，返回 markdown 文本；未装/推理失败返回 None（调用方 fallback fitz）。

    出处：2026-07-01 MinerU 因 `requires_python<3.14` 跟本仓库 Python 3.14 环境硬不兼容（`pip install`
    直接被拒，非网络/代理问题），找替代时 pymupdf4llm 的增强路径又卡在硬依赖 `pymupdf_layout` 是
    Polyform Noncommercial 许可（免费只能非商业用，同 D6 因 GPL-3.0 拒 AiPPT 是同一类法务风险，商用
    需买 Artifex 商业授权）——docling 是三个候选里唯一 Python 3.14 装得上 + MIT 许可干净的。
    **真实验证过**（非只凭 license 幸存者偏差二选一）：合成 PDF 用 PyMuPDF 内置简化字体测出过乱码，
    换真实系统字体（`STHeiti Light.ttc`）+ 仓库里真实 38 页 McKinsey 咨询 deck 重测，中文正确提取、
    从真实文档里拆出 118 行结构化表格（含一张 13 列月度财务假设表），证明复杂版面/表格能力是真的。
    详见 `specs/决策记录.md` D13。⚠️ 代价：依赖链重（torch 单独 443MB，全装完 1GB+），比 MinerU 本身
    还重，这不是"零成本平替"，是"扛得住重量换来能用+许可证干净"的权衡。
    """
    try:
        from docling.document_converter import DocumentConverter
    except ImportError:
        return None  # 未装 docling 包
    try:
        result = DocumentConverter().convert(str(path))
        return result.document.export_to_markdown()
    except Exception:
        return None  # 模型权重未下载(需联网)/推理失败等——不让 brief 解析被 docling 内部错误拖垮


def _parse_pdf_with_mineru(path: Path) -> str | None:
    """真调用 MinerU pipeline 后端解析 PDF，返回 markdown 文本；未装/推理失败返回 None（调用方 fallback fitz）。

    出处【mineru/cli/common.py::do_parse】：CLI 导向的批处理 API——写文件到 output_dir，不直接回传文本，
    需要读回 `{stem}/{parse_method}/{stem}.md`（出处【mineru/cli/common.py::prepare_env】的目录约定）。
    用 `rglob("*.md")` 找回文件而非硬编码精确文件名——`do_parse` 内部对文件名做过长截断/去重，
    精确 stem 不保证跟输入文件名一致，扫目录比赌文件名更稳。

    backend="pipeline"：唯一纯 CPU 可跑、"无幻觉"的后端（vlm/hybrid 需 8GB+ 显存）——契合 brief 解析
    "宁可漏识别不可瞎编"的取向，见 07 文档预处理/MinerU 深拆笔记。
    """
    try:
        from mineru.cli.common import do_parse, read_fn
    except ImportError:
        return None  # 未装 mineru 包
    import tempfile
    with tempfile.TemporaryDirectory() as tmp:
        try:
            pdf_bytes = read_fn(str(path))
            do_parse(
                output_dir=tmp,
                pdf_file_names=[path.stem],
                pdf_bytes_list=[pdf_bytes],
                p_lang_list=["ch"],
                backend="pipeline",
                f_draw_layout_bbox=False, f_draw_span_bbox=False,
                f_dump_middle_json=False, f_dump_model_output=False,
                f_dump_orig_pdf=False, f_dump_content_list=False,
            )
        except Exception:
            return None  # 模型权重未下载/推理失败等——不让 brief 解析被 MinerU 内部错误拖垮
        md_files = list(Path(tmp).rglob("*.md"))
        if not md_files:
            return None
        return md_files[0].read_text(encoding="utf-8")


def parse_brief(path, *, prefer_mineru: bool = False, prefer_docling: bool = False) -> dict:
    """brief 文件 → {text, format, pages}。按扩展名选提取器（免费本地·无 key）。
    支持 pdf/pptx/docx/txt/md；老二进制 .ppt 明确拒收（python-pptx 打不开·提示先另存为 .pptx）。

    prefer_docling=True 时 PDF 优先走 docling 版面分析增强（复杂版面/表格更准·`pip install docling`·
    MIT 许可·2026-07-01 起**推荐路径**，真实验证过中文+复杂表格，见 D13 决策记录）。
    prefer_mineru=True 时走 MinerU（`pip install mineru`——⚠️ 本仓库 Python 3.14 环境下 mineru 包
    本身装不上，此参数留着是因为其他 Python<3.14 环境仍可用，不是本仓库当前可行路径）。
    两者都未装/失败则 fallback fitz（轻量·已实战够用）。同时置 True 时 docling 优先（更新推荐）。
    """
    path = Path(path)
    ext = path.suffix.lower()

    if ext == ".pdf":
        if prefer_docling:
            docling_text = _parse_pdf_with_docling(path)
            if docling_text is not None:
                return {"text": docling_text, "format": "pdf", "pages": None, "parser": "docling"}
        if prefer_mineru:
            mineru_text = _parse_pdf_with_mineru(path)
            if mineru_text is not None:
                return {"text": mineru_text, "format": "pdf", "pages": None, "parser": "mineru"}
        import fitz
        d = fitz.open(str(path))
        pages = [p.get_text() for p in d]
        return {"text": "\n\n".join(pages), "format": "pdf", "pages": len(pages), "parser": "fitz"}

    if ext == ".ppt":
        # 2026-07-02 saopan扫盘揪出：旧版把 .ppt 收进 .pptx 分支——python-pptx 只认 OOXML（zip 容器），
        # 老二进制 .ppt（OLE 复合文档）一喂必崩，且崩的还是库内部的 zipfile 报错，用户看不懂。
        # fail-closed：明确拒收 + 给出用户能照做的出路。
        raise ValueError("老二进制 .ppt 不支持（python-pptx 只认 OOXML 格式）·"
                         "请先在 PowerPoint 里「另存为 .pptx」再传入")

    if ext == ".pptx":
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(str(path))
        chunks = []
        for i, s in enumerate(prs.slides, 1):
            parts = []
            n_pics = 0
            for sh in _iter_pptx_shapes(s.shapes):  # 递归展开 GroupShape·组内文字/图片不再静默丢
                try:
                    if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        n_pics += 1
                        continue
                    text = _pptx_shape_text(sh)
                    if text:
                        parts.append(text)
                except Exception:
                    continue  # python-pptx 3.14 个别形状读回兼容·跳过不中断
            if n_pics:
                parts.append(f"[本页另有 {n_pics} 张图片·brief 提取不含图片内容，只记存在]")
            try:
                if s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip():
                    parts.append(f"[演讲者备注] {s.notes_slide.notes_text_frame.text}")
            except Exception:
                pass
            if parts:
                chunks.append(f"[slide {i}]\n" + "\n".join(parts))
        return {"text": "\n\n".join(chunks), "format": "pptx", "pages": len(prs.slides)}

    if ext == ".docx":
        import docx
        doc = docx.Document(str(path))
        parts = []
        for p in doc.paragraphs:
            text = _docx_para_text(p)
            if text.strip():
                parts.append(text)
        for t in doc.tables:
            md = _docx_table_md(t)
            if md:
                parts.append(md)
        n_imgs = len(doc.inline_shapes)
        if n_imgs:
            parts.append(f"[文档另有 {n_imgs} 张图片·brief 提取不含图片内容，只记存在]")
        return {"text": "\n".join(parts), "format": "docx", "pages": None}

    if ext in (".txt", ".md"):
        return {"text": path.read_text(encoding="utf-8"), "format": "text", "pages": None}

    raise ValueError(f"不支持的 brief 格式 {ext}（支持 pdf/pptx/docx/txt/md）")
