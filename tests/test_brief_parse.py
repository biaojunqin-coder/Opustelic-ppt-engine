"""brief 解析 + 调研 provider（策略阶段1/3·免费方案）测试。"""

from __future__ import annotations

import sys
import types
from pathlib import Path

import pytest

from engine.brief_parse import parse_brief
from engine.research import attach_citation, attach_source, research

try:
    import docling.document_converter  # noqa: F401
    _HAS_DOCLING = True
except ImportError:
    _HAS_DOCLING = False

_CJK_SYSTEM_FONT = "/System/Library/Fonts/STHeiti Light.ttc"


def test_parse_txt(tmp_path):
    f = tmp_path / "brief.txt"
    f.write_text("CFO 要做成本管控 deck", encoding="utf-8")
    r = parse_brief(f)
    assert r["format"] == "text" and "成本管控" in r["text"]


def test_parse_md(tmp_path):
    f = tmp_path / "b.md"
    f.write_text("# brief\n要做融资 deck", encoding="utf-8")
    assert "融资" in parse_brief(f)["text"]


def _make_pdf(tmp_path, text="这是brief正文内容·数字42"):
    import fitz
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), text)
    f = tmp_path / "brief.pdf"
    doc.save(str(f))
    doc.close()
    return f


def test_parse_pdf_default_uses_fitz(tmp_path):
    r = parse_brief(_make_pdf(tmp_path))
    assert r["format"] == "pdf" and r["parser"] == "fitz" and "42" in r["text"] and r["pages"] == 1


def _install_fake_mineru(monkeypatch, *, raise_on_parse=False, md_content="# mineru解析结果·表格更准"):
    """向 sys.modules 注入假 mineru.cli.common，模拟"已装 mineru"时 do_parse 成功/失败两条路径——
    真实模型推理需要下载的权重不在本仓库开发环境里，这里只验证 `_parse_pdf_with_mineru` 的接线
    （调用参数对不对/读回文件对不对/异常时正确 fallback），不是端到端真推理。"""
    from pathlib import Path

    fake_common = types.ModuleType("mineru.cli.common")

    def fake_read_fn(path):
        return b"fake-pdf-bytes"

    def fake_do_parse(*, output_dir, pdf_file_names, pdf_bytes_list, p_lang_list, backend, **kwargs):
        assert backend == "pipeline"  # 必须用纯CPU可跑+无幻觉的后端，不是vlm/hybrid
        if raise_on_parse:
            raise RuntimeError("模拟模型权重未下载/推理失败")
        stem = pdf_file_names[0]
        out_dir = Path(output_dir) / stem / "auto"
        out_dir.mkdir(parents=True, exist_ok=True)
        (out_dir / f"{stem}.md").write_text(md_content, encoding="utf-8")

    fake_common.do_parse = fake_do_parse
    fake_common.read_fn = fake_read_fn
    fake_cli = types.ModuleType("mineru.cli")
    fake_cli.common = fake_common
    fake_mineru = types.ModuleType("mineru")
    fake_mineru.cli = fake_cli

    monkeypatch.setitem(sys.modules, "mineru", fake_mineru)
    monkeypatch.setitem(sys.modules, "mineru.cli", fake_cli)
    monkeypatch.setitem(sys.modules, "mineru.cli.common", fake_common)


def test_parse_pdf_prefer_mineru_not_installed_falls_back_to_fitz(tmp_path):
    # 本仓库开发环境未装 mineru 包·真实 ImportError 路径（非 mock）
    r = parse_brief(_make_pdf(tmp_path), prefer_mineru=True)
    assert r["parser"] == "fitz" and "42" in r["text"]


def test_parse_pdf_prefer_mineru_success_uses_mineru_output(tmp_path, monkeypatch):
    _install_fake_mineru(monkeypatch, md_content="# mineru解析结果·表格更准")
    r = parse_brief(_make_pdf(tmp_path), prefer_mineru=True)
    assert r["parser"] == "mineru" and "mineru解析结果" in r["text"] and r["pages"] is None


def test_parse_pdf_prefer_mineru_runtime_failure_falls_back_to_fitz(tmp_path, monkeypatch):
    # 模拟"装了 mineru 但模型权重未下载/推理报错"——不能让 brief 解析被内部错误拖垮，必须 fallback
    _install_fake_mineru(monkeypatch, raise_on_parse=True)
    r = parse_brief(_make_pdf(tmp_path), prefer_mineru=True)
    assert r["parser"] == "fitz" and "42" in r["text"]


# ── docling（2026-07-01·D13 决策记录·MinerU 因 Python3.14 硬不兼容后找到的替代增强路径）──

def _install_fake_docling(monkeypatch, *, raise_on_convert=False, md_content="# docling解析结果·表格更准"):
    """向 sys.modules 注入假 docling.document_converter，模拟"已装 docling"时 convert 成功/失败两条
    路径——不依赖本机是否真装了 docling 包，跟 mineru 的 mock 测试同一个目的：只验证
    `_parse_pdf_with_docling` 的接线（调用对不对/异常时正确 fallback），不是端到端真推理
    （真推理见下面 `test_parse_pdf_prefer_docling_real_integration_extracts_chinese`）。"""

    class FakeDocument:
        def export_to_markdown(self):
            return md_content

    class FakeResult:
        document = FakeDocument()

    class FakeConverter:
        def convert(self, path):
            if raise_on_convert:
                raise RuntimeError("模拟模型权重未下载/推理失败")
            return FakeResult()

    fake_module = types.ModuleType("docling.document_converter")
    fake_module.DocumentConverter = FakeConverter
    fake_docling = types.ModuleType("docling")
    fake_docling.document_converter = fake_module

    monkeypatch.setitem(sys.modules, "docling", fake_docling)
    monkeypatch.setitem(sys.modules, "docling.document_converter", fake_module)


def test_parse_pdf_prefer_docling_not_installed_falls_back_to_fitz(tmp_path, monkeypatch):
    # sys.modules[name]=None 强制模拟"未装"——本机可能真装了 docling，用这个确保测试不依赖本机状态
    monkeypatch.setitem(sys.modules, "docling", None)
    monkeypatch.setitem(sys.modules, "docling.document_converter", None)
    r = parse_brief(_make_pdf(tmp_path), prefer_docling=True)
    assert r["parser"] == "fitz" and "42" in r["text"]


def test_parse_pdf_prefer_docling_success_uses_docling_output(tmp_path, monkeypatch):
    _install_fake_docling(monkeypatch, md_content="# docling解析结果·表格更准")
    r = parse_brief(_make_pdf(tmp_path), prefer_docling=True)
    assert r["parser"] == "docling" and "docling解析结果" in r["text"] and r["pages"] is None


def test_parse_pdf_prefer_docling_runtime_failure_falls_back_to_fitz(tmp_path, monkeypatch):
    _install_fake_docling(monkeypatch, raise_on_convert=True)
    r = parse_brief(_make_pdf(tmp_path), prefer_docling=True)
    assert r["parser"] == "fitz" and "42" in r["text"]


def test_parse_pdf_prefer_docling_beats_mineru_when_both_true(tmp_path, monkeypatch):
    # 两者都开时 docling 优先——2026-07-01起是更新推荐的路径（D13 决策记录）
    _install_fake_docling(monkeypatch, md_content="# docling输出")
    _install_fake_mineru(monkeypatch, md_content="# mineru输出")
    r = parse_brief(_make_pdf(tmp_path), prefer_docling=True, prefer_mineru=True)
    assert r["parser"] == "docling"


@pytest.mark.skipif(not _HAS_DOCLING, reason="docling 未装·真实解析验证跳过")
@pytest.mark.skipif(not Path(_CJK_SYSTEM_FONT).exists(), reason="非macOS环境·系统CJK字体路径不存在")
def test_parse_pdf_prefer_docling_real_integration_extracts_chinese(tmp_path):
    """真实端到端验证(非mock)：用真实系统CJK字体(非PyMuPDF内置简化字体)造PDF。
    2026-07-01 踩过两个坑，都记下来——① PyMuPDF内置简化字体(fontname="china-ss")的内部编码方式
    会让 docling 解析出乱码，换真实系统字体(Word/PPT导出实际会用的那类)才对；② 单行文字贴在页面
    顶部会被docling版面分类器正确识别成"page_header"(页眉)从而被`export_to_markdown()`过滤掉——
    这不是bug是合理设计(页眉通常非正文)，测试fixture要放在页面中段模拟真实正文位置才对。
    仓库里真实38页McKinsey咨询deck另测过中文和复杂表格(118行结构化表格)都正确提取，
    见 D13 决策记录完整验证过程。"""
    import fitz
    doc = fitz.open()
    page = doc.new_page()
    page.insert_font(fontname="F0", fontfile=_CJK_SYSTEM_FONT)
    page.insert_text((72, 300), "真实文档中文测试·收入增速25%，需要专项管控。", fontsize=14, fontname="F0")
    f = tmp_path / "cjk.pdf"
    doc.save(str(f))
    r = parse_brief(f, prefer_docling=True)
    assert r["parser"] == "docling" and "真实文档中文测试" in r["text"] and "25%" in r["text"]


def test_parse_unsupported(tmp_path):
    f = tmp_path / "x.xyz"; f.write_text("x")
    with pytest.raises(ValueError, match="不支持"):
        parse_brief(f)


# ── B2-10 回归（2026-07-02 saopan扫盘）：旧版把 .ppt 收进 .pptx 分支——python-pptx 只认 OOXML，
#    老二进制 .ppt 一喂必崩且崩的是 zipfile 内部报错；现在明确拒收 + 给出用户能照做的出路 ──
def test_parse_ppt_rejected_with_actionable_message(tmp_path):
    f = tmp_path / "old.ppt"
    f.write_bytes(b"\xd0\xcf\x11\xe0legacy-ole-binary")  # OLE 复合文档魔数·真实 .ppt 的开头
    with pytest.raises(ValueError, match="另存为 .pptx"):
        parse_brief(f)


# ── B2-11 回归（2026-07-02 saopan扫盘）：GroupShape 没有 text_frame，组合形状里的文字旧版整组
#    静默丢失（真实 brief 里"几个框圈成一组"极常见）——现在递归展开组（含嵌套组）──
def _pptx_with_grouped_text(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tx = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    tx.text_frame.text = "顶层文字"

    group = slide.shapes.add_group_shape()
    gtx = group.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(1))
    gtx.text_frame.text = "组内文字X"
    inner = group.shapes.add_group_shape()  # 嵌套组·递归也要进得去
    itx = inner.shapes.add_textbox(Inches(1), Inches(3), Inches(3), Inches(1))
    itx.text_frame.text = "嵌套组内文字Y"

    f = tmp_path / "grouped.pptx"
    prs.save(f)
    return f


def test_parse_pptx_group_shape_text_extracted(tmp_path):
    r = parse_brief(_pptx_with_grouped_text(tmp_path))
    assert "顶层文字" in r["text"]
    assert "组内文字X" in r["text"]        # 旧版这行静默丢失
    assert "嵌套组内文字Y" in r["text"]    # 组中组也要提出来


# ── 表格/图片清单/超链接（2026-07-01·参考 ppt-master source_to_md 补的缺口·之前静默丢失）──
def _pptx_with_table_image_link_notes(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches
    from PIL import Image

    img = tmp_path / "tiny.png"
    Image.new("RGB", (4, 4), color="red").save(img)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tx = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    p = tx.text_frame.paragraphs[0]
    p.add_run().text = "前文 "
    link_run = p.add_run(); link_run.text = "链接"
    link_run.hyperlink.address = "https://example.com/y"

    table = slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(3), Inches(1)).table
    table.cell(0, 0).text = "指标"; table.cell(0, 1).text = "数值"
    table.cell(1, 0).text = "收入"; table.cell(1, 1).text = "100万"

    slide.shapes.add_picture(str(img), Inches(1), Inches(3))
    slide.notes_slide.notes_text_frame.text = "这是演讲者备注内容"

    f = tmp_path / "full.pptx"
    prs.save(f)
    return f


def test_parse_pptx_table_extracted_as_markdown(tmp_path):
    r = parse_brief(_pptx_with_table_image_link_notes(tmp_path))
    assert "| 指标 | 数值 |" in r["text"] and "| 收入 | 100万 |" in r["text"]


def test_parse_pptx_hyperlink_inline(tmp_path):
    r = parse_brief(_pptx_with_table_image_link_notes(tmp_path))
    assert "[链接](https://example.com/y)" in r["text"]


def test_parse_pptx_image_noted_not_silently_dropped(tmp_path):
    r = parse_brief(_pptx_with_table_image_link_notes(tmp_path))
    assert "1 张图片" in r["text"]


def test_parse_pptx_speaker_notes_included(tmp_path):
    r = parse_brief(_pptx_with_table_image_link_notes(tmp_path))
    assert "这是演讲者备注内容" in r["text"]


def _add_docx_hyperlink(paragraph, text, url):
    """python-docx 无原生写超链接 API，按其内部 OOXML 结构手动拼（只用于造测试夹具，读路径走公开 API）。"""
    from docx.opc.constants import RELATIONSHIP_TYPE as RT
    from docx.oxml.ns import qn
    r_id = paragraph.part.relate_to(url, RT.HYPERLINK, is_external=True)
    hyperlink = paragraph._p.makeelement(qn("w:hyperlink"), {qn("r:id"): r_id})
    run = paragraph._p.makeelement(qn("w:r"), {})
    t = paragraph._p.makeelement(qn("w:t"), {})
    t.text = text
    run.append(t)
    hyperlink.append(run)
    paragraph._p.append(hyperlink)


def _docx_with_table_image_link(tmp_path):
    import docx
    from docx.shared import Inches
    from PIL import Image

    img = tmp_path / "tiny.png"
    Image.new("RGB", (4, 4), color="blue").save(img)

    doc = docx.Document()
    p = doc.add_paragraph("前文 ")
    _add_docx_hyperlink(p, "点这里", "https://example.com/x")
    p.add_run(" 后文")

    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "项目"; table.cell(0, 1).text = "状态"
    table.cell(1, 0).text = "立项"; table.cell(1, 1).text = "已完成"

    doc.add_picture(str(img), width=Inches(1))

    f = tmp_path / "b.docx"
    doc.save(f)
    return f


def test_parse_docx_table_extracted_as_markdown(tmp_path):
    r = parse_brief(_docx_with_table_image_link(tmp_path))
    assert "| 项目 | 状态 |" in r["text"] and "| 立项 | 已完成 |" in r["text"]


def test_parse_docx_hyperlink_inline(tmp_path):
    r = parse_brief(_docx_with_table_image_link(tmp_path))
    assert "前文 [点这里](https://example.com/x) 后文" in r["text"]


def test_parse_docx_image_noted_not_silently_dropped(tmp_path):
    r = parse_brief(_docx_with_table_image_link(tmp_path))
    assert "1 张图片" in r["text"]


def test_research_stub():
    assert research("人力成本趋势")["status"].startswith("stub")


def test_research_claude_websearch():
    assert "claude-websearch" in research("人力成本趋势", provider="claude-websearch")["status"]


def test_research_unknown_provider_raises():
    with pytest.raises(NotImplementedError):
        research("q", provider="gpt-researcher")


def test_attach_source():
    assert attach_source({"dim": "净利", "data": "-20%"}, "财报2024")["source"] == "财报2024"


def test_research_today_anchors_note():
    # today 真实生效(非装饰性参数)：写进返回 note，提醒下游"某年预计"挂过期的坑
    r = research("市场规模", provider="claude-websearch", today="2026-06-30")
    assert r["today"] == "2026-06-30" and "2026-06-30" in r["note"]


def test_research_no_today_warns_in_note():
    r = research("市场规模", provider="claude-websearch")
    assert "未传 today" in r["note"]


def test_attach_citation_full_fields():
    ev = attach_citation({"dim": "游客", "data": "+21.57%"}, title="标题", outlet="文旅部官网",
                         date="2026-03-10", quote="原文片段", url="mct.gov.cn/x")
    assert ev["source_title"] == "标题" and ev["source_url"] == "mct.gov.cn/x" and ev["source"] == "文旅部官网"
